"""Generate the OT Relay (電驛) Protection PRD deck (zh-TW, dark SenseL EdgeX theme).

Field-oriented PRD: solution overview + deployment topology + configuration.
Run: python3 docs/build_relay_prd_deck.py
Output: docs/SenseL-EdgeX-Relay-Protection-PRD.pptx
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "assets", "sensel-edgex-logo.png")
OUT = os.path.join(HERE, "SenseL-EdgeX-Relay-Protection-PRD.pptx")

EMU = 914400
SW, SH = int(13.333 * EMU), int(7.5 * EMU)

BG = RGBColor(0x0E, 0x16, 0x21)
BG2 = RGBColor(0x16, 0x20, 0x2E)
CARD = RGBColor(0x1A, 0x25, 0x36)
GREEN = RGBColor(0x9B, 0xD5, 0x34)
PURPLE = RGBColor(0xA2, 0x4B, 0xD0)
BLUE = RGBColor(0x4D, 0x9B, 0xF0)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x5D, 0x5D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9C, 0xA8, 0xB8)
LINE = RGBColor(0x2A, 0x36, 0x48)

EA = "PingFang TC"   # east-asian typeface
LAT = "Aptos"

prs = Presentation()
prs.slide_width = Emu(SW)
prs.slide_height = Emu(SH)
BLANK = prs.slide_layouts[6]


def E(v):
    return Emu(int(v * EMU))


def _ea(run, font):
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", font)


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(SW), Emu(SH))
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, round_=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(st, E(x), E(y), E(w), E(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sa=3):
    tb = s.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0)
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = LAT; _ea(r, EA)
    return tb


def logo_corner(s):
    w = 1.85; h = w * 682 / 1024
    s.shapes.add_picture(LOGO, E(13.333 - w - 0.3), E(0.28), E(w), E(h))


def header(s, kicker, title):
    rect(s, 0.55, 0.6, 0.12, 0.92, fill=GREEN)
    txt(s, 0.85, 0.54, 9.5, 0.4, [[(kicker, 13, GREEN, True)]])
    txt(s, 0.83, 0.9, 10.2, 0.75, [[(title, 27, WHITE, True)]])
    logo_corner(s)
    rect(s, 0.85, 1.62, 11.6, 0.016, fill=LINE)


def footer(s, idx):
    txt(s, 0.55, 7.04, 9.0, 0.32,
        [[("SenseL EdgeX (RelayGuard) ·  OT 電驛設備防護 PRD ·  機密", 9, MUTED, False)]])
    txt(s, 12.2, 7.04, 0.6, 0.32, [[(f"{idx:02d}", 10, MUTED, True)]], align=PP_ALIGN.RIGHT)


def banner(s, y, runs, color=GREEN):
    rect(s, 0.85, y, 11.6, 0.6, fill=CARD, line=color, lw=1.0, round_=True)
    txt(s, 1.1, y, 11.1, 0.6, [runs], anchor=MSO_ANCHOR.MIDDLE)


def card(s, x, y, w, h, title, body, color, num=None):
    rect(s, x, y, w, h, fill=CARD, line=LINE, lw=0.75, round_=True)
    tx = x + 0.22
    if num is not None:
        b = rect(s, x + 0.2, y + 0.2, 0.42, 0.42, fill=color)
        tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(num); r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = BG; r.font.name = LAT
        tx = x + 0.74
    txt(s, tx, y + 0.2, w - (tx - x) - 0.2, 0.4, [[(title, 14.5, color, True)]])
    txt(s, x + 0.22, y + 0.66, w - 0.44, h - 0.8, [[(body, 11.5, MUTED, False)]])


def chip(s, x, y, w, h, lines, color):
    rect(s, x, y, w, h, fill=BG2, line=color, lw=1.25, round_=True)
    runs = [[(t, sz, c, b)] for (t, sz, c, b) in lines]
    txt(s, x + 0.05, y, w - 0.1, h, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sa=1)


# ═══════════ Slide 1 — 封面 ═══════════
s = slide()
rect(s, 0, 0, 13.333, 0.10, fill=GREEN)
rect(s, 0, 7.40, 13.333, 0.10, fill=PURPLE)
lw_ = 4.6; lh_ = lw_ * 682 / 1024
s.shapes.add_picture(LOGO, E((13.333 - lw_) / 2), E(1.05), E(lw_), E(lh_))
txt(s, 1.0, 4.0, 11.333, 0.8, [[("OT 電驛設備防護解決方案", 32, WHITE, True)]], align=PP_ALIGN.CENTER)
txt(s, 1.0, 4.85, 11.333, 0.5,
    [[("產品需求文件 (PRD) ·  場域解決方案與配置說明", 16, GREEN, True)]], align=PP_ALIGN.CENTER)
txt(s, 1.0, 5.7, 11.333, 0.4,
    [[("被動式邊緣偵測 ·  IEC 61850 / Modbus ·  變電所與工控場域", 13, MUTED, False)]], align=PP_ALIGN.CENTER)
txt(s, 1.0, 6.5, 11.333, 0.35, [[("Draft v0.1 ·  by AvocadoAI ·  機密", 11, MUTED, False)]], align=PP_ALIGN.CENTER)

# ═══════════ Slide 2 — 背景與問題 ═══════════
s = slide()
header(s, "背景與問題", "電驛 (Protective Relay / IED) 為什麼需要被動防護")
cards = [
    ("攻擊面擴大", "變電所數位化後，保護電驛 (IED) 透過 IEC 61850 GOOSE/MMS、Modbus 連網，暴露於橫向移動與惡意指令風險。", RED),
    ("傳統設備無防護", "電驛多為封閉韌體、無法安裝代理程式 (agentless)，且不可承受主動掃描，傳統 IT 資安工具無法套用。", AMBER),
    ("缺乏可視性", "現場缺乏被動監控，無法得知誰在存取電驛、是否有未授權寫入、GOOSE 是否異常或 IED 是否離線。", BLUE),
    ("無法中斷生產", "電驛關乎供電與人身安全，任何防護方案必須『零干擾、不阻斷』，僅能被動觀測與告警。", PURPLE),
]
x = 0.85
for t, b, c in cards[:2]:
    card(s, x, 1.85, 5.65, 1.85, t, b, c)
    x += 5.95
x = 0.85
for t, b, c in cards[2:]:
    card(s, x, 3.85, 5.65, 1.85, t, b, c)
    x += 5.95
banner(s, 5.95, ("核心需求：在不接觸、不掃描、不阻斷電驛的前提下，取得可視性、偵測異常並留存證據。", 14, WHITE, True))
footer(s, 2)

# ═══════════ Slide 3 — 產品定位與範圍 ═══════════
s = slide()
header(s, "產品定位", "RelayGuard：被動式電驛防護邊緣閘道")
txt(s, 0.85, 1.78, 11.6, 0.6,
    [[("部署於變電所 / 廠區的邊緣閘道，透過 SPAN/TAP 鏡像被動解析電驛通訊，結合 EdgeX 主動遙測，"
       "提供資產可視性、異常偵測、威脅情資比對與證據留存。", 14, MUTED, False)]])
# scope vs non-scope
rect(s, 0.85, 2.65, 5.65, 3.5, fill=CARD, line=GREEN, lw=1.0, round_=True)
txt(s, 1.1, 2.8, 5.2, 0.4, [[("MVP 範圍", 16, GREEN, True)]])
for i, t in enumerate([
    "SPAN/TAP 被動鏡像解析 (L2–L7)",
    "IEC 61850 (GOOSE/MMS)、Modbus TCP",
    "OT-001~019 規則偵測與告警",
    "資產被動發現 + EdgeX 設備納管",
    "防護基線 (baseline) 與政策下派",
    "CTI 威脅情資比對 (OT-019)",
    "PCAP 證據留存、事件 / 健康上報",
]):
    txt(s, 1.15, 3.25 + i * 0.4, 5.1, 0.38, [[("✓  ", 12, GREEN, True), (t, 12, WHITE, False)]])
rect(s, 6.8, 2.65, 5.65, 3.5, fill=CARD, line=AMBER, lw=1.0, round_=True)
txt(s, 7.05, 2.8, 5.2, 0.4, [[("MVP 非目標 (後續階段)", 16, AMBER, True)]])
for i, t in enumerate([
    "自動阻斷 / 隔離 (後續結合 ZTNA)",
    "1 Gbps 線速 DPI",
    "加密 OPC-UA payload 被動解密",
    "完整 IEC 61850 SV / SCL 語意 decoder",
    "HA / 多站點叢集",
]):
    txt(s, 7.1, 3.25 + i * 0.45, 5.1, 0.4, [[("•  ", 12, AMBER, True), (t, 12, MUTED, False)]])
footer(s, 3)

# ═══════════ Slide 4 — 解決方案總覽 (雙路徑) ═══════════
s = slide()
header(s, "解決方案總覽", "雙路徑架構：被動鏡像偵測 + 主動遙測")
# field device
chip(s, 0.85, 2.05, 2.1, 1.1, [("電驛 / IED", 14, WHITE, True), ("PLC · RTU · HMI", 10, MUTED, False)], BLUE)
# passive path
a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, E(3.0), E(2.3), E(0.55), E(0.22))
a.fill.solid(); a.fill.fore_color.rgb = GREEN; a.line.fill.background(); a.shadow.inherit = False
chip(s, 3.65, 1.65, 2.4, 0.85, [("SPAN / TAP", 12, GREEN, True), ("鏡像唯讀擷取", 9, MUTED, False)], GREEN)
chip(s, 3.65, 2.7, 2.4, 0.85, [("EdgeX 主動輪詢", 12, PURPLE, True), ("Modbus / OPC-UA", 9, MUTED, False)], PURPLE)
a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, E(6.1), E(2.3), E(0.45), E(0.22))
a.fill.solid(); a.fill.fore_color.rgb = MUTED; a.line.fill.background(); a.shadow.inherit = False
# edge box
rect(s, 6.6, 1.5, 3.4, 2.25, fill=CARD, line=GREEN, lw=1.25, round_=True)
txt(s, 6.8, 1.6, 3.0, 0.35, [[("RelayGuard 邊緣閘道", 13, GREEN, True)]])
for i, t in enumerate(["Packet Sensor 解析 + 偵測", "EdgeX 遙測正規化", "防護基線 / 政策套用", "PCAP 證據 ring buffer"]):
    txt(s, 6.85, 2.0 + i * 0.4, 3.0, 0.38, [[("• ", 11, GREEN, True), (t, 11, WHITE, False)]])
a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, E(10.05), E(2.5), E(0.45), E(0.22))
a.fill.solid(); a.fill.fore_color.rgb = MUTED; a.line.fill.background(); a.shadow.inherit = False
chip(s, 10.55, 1.9, 1.9, 1.4, [("Control", 13, WHITE, True), ("Plane /", 13, WHITE, True), ("Portal", 13, WHITE, True)], BLUE)
banner(s, 4.25, ("設計原則：原始鏡像流量不進 EdgeX；僅上傳解析後的特徵摘要、安全事件與證據 metadata。", 13.5, WHITE, True))
txt(s, 0.85, 5.15, 11.6, 1.0,
    [[("上報路徑：", 13, GREEN, True),
      ("安全事件與遙測以 MQTT (北向) 為主、HTTP 為輔上傳 Control Plane；離線時本地 SQLite 緩衝，恢復後補送。"
       "下行則接收 CTI 黑名單與偵測政策。", 13, MUTED, False)]])
footer(s, 4)

# ═══════════ Slide 5 — 電驛防護偵測能力 (規則表) ═══════════
s = slide()
header(s, "偵測能力", "電驛相關偵測規則 (OT Detection Rules)")
rows = [
    ("OT-007", "未預期 Modbus 寫入", "high", "對電驛寫入非授權功能碼 / 暫存器"),
    ("OT-009", "電驛離線 (Relay offline)", "high", "遙測 + 被動觀測判定設備離線"),
    ("OT-010", "未授權主機存取電驛", "high", "非白名單來源連線至電驛"),
    ("OT-012", "GOOSE test bit 出現於生產", "high", "正式環境出現測試旗標"),
    ("OT-016", "未預期 MMS 寫入", "high", "對 IED 寫入控制 / 設定值"),
    ("OT-017", "GOOSE 靜默 (IED offline)", "high", "GOOSE 超過 max_silence 未發布"),
    ("OT-018", "未授權 MMS 連線至電驛 IED", "high", "非白名單 MMS client 連線"),
    ("OT-019", "CTI IoC 命中", "high", "比對情資黑名單 (IP/domain/hash)"),
]
# header
hx = [0.85, 2.1, 6.2, 7.7]; hw = [1.25, 4.1, 1.5, 4.75]
for i, h in enumerate(["Rule", "名稱", "嚴重度", "觸發情境"]):
    rect(s, hx[i], 1.75, hw[i], 0.42, fill=PURPLE)
    txt(s, hx[i] + 0.1, 1.75, hw[i] - 0.2, 0.42, [[(h, 11.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
y = 2.17
for rid, name, sev, sc in rows:
    rect(s, 0.85, y, 11.6, 0.46, fill=CARD if (int(rid[3:]) % 2 == 0) else BG2, line=LINE, lw=0.5)
    txt(s, 0.95, y, 1.15, 0.46, [[(rid, 11, GREEN, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 2.2, y, 4.0, 0.46, [[(name, 11.5, WHITE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 6.3, y, 1.4, 0.46, [[(sev.upper(), 10.5, RED if sev == "high" else AMBER, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 7.8, y, 4.6, 0.46, [[(sc, 11, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.5
txt(s, 0.85, y + 0.05, 11.6, 0.5,
    [[("另含 OT-001~006 / OT-008 / OT-011 / OT-013~015 等網路行為與 IEC 61850 規則；完整 19 條規則於產品文件。",
       10.5, MUTED, False)]])
footer(s, 5)

# ═══════════ Slide 6 — IEC 61850 專項偵測 ═══════════
s = slide()
header(s, "IEC 61850 專項", "GOOSE / MMS 電驛通訊深度偵測")
rect(s, 0.85, 1.85, 5.65, 4.1, fill=CARD, line=GREEN, lw=1.0, round_=True)
txt(s, 1.1, 2.0, 5.2, 0.4, [[("GOOSE (L2 多點傳播)", 16, GREEN, True)]])
for i, (t, d) in enumerate([
    ("OT-011 新 GOOSE publisher", "出現未登記的發布者 MAC/APPID"),
    ("OT-012 test bit 於生產", "正式環境出現測試旗標 (高風險)"),
    ("OT-013 stNum 異常跳動", "狀態序號超過門檻，疑似偽造"),
    ("OT-017 GOOSE 靜默", "超過 max_silence_sec → IED 離線"),
]):
    txt(s, 1.15, 2.5 + i * 0.82, 5.1, 0.4, [[(t, 12.5, WHITE, True)]])
    txt(s, 1.15, 2.82 + i * 0.82, 5.1, 0.42, [[(d, 11, MUTED, False)]])
rect(s, 6.8, 1.85, 5.65, 4.1, fill=CARD, line=PURPLE, lw=1.0, round_=True)
txt(s, 7.05, 2.0, 5.2, 0.4, [[("MMS (TCP/102 主從)", 16, PURPLE, True)]])
for i, (t, d) in enumerate([
    ("OT-014 新 MMS client", "未登記的 client 連線至 IED"),
    ("OT-015 MMS session 異常", "連線速率超過門檻"),
    ("OT-016 未預期 MMS 寫入", "對 IED 寫入控制 / 設定值 (高風險)"),
    ("OT-018 未授權 MMS 至電驛", "非白名單 client 存取保護 IED"),
]):
    txt(s, 7.1, 2.5 + i * 0.82, 5.1, 0.4, [[(t, 12.5, WHITE, True)]])
    txt(s, 7.1, 2.82 + i * 0.82, 5.1, 0.42, [[(d, 11, MUTED, False)]])
banner(s, 6.15, ("以每台 IED 的白名單 (publisher MAC/APPID、允許 MMS client) 與門檻為基準，"
                 "偏離即告警；不對 IED 發出任何主動封包。", 13, WHITE, True))
footer(s, 6)

# ═══════════ Slide 7 — 場域部署拓撲 ═══════════
s = slide()
header(s, "場域部署", "變電所 / 廠區網路接線拓撲")
# OT switch
chip(s, 1.4, 2.0, 2.5, 1.0, [("OT 交換器", 13, WHITE, True), ("變電所 / 機房", 10, MUTED, False)], BLUE)
# devices below
for i, name in enumerate(["保護電驛 IED", "Modbus RTU", "HMI / 工作站"]):
    chip(s, 0.85 + i * 1.05, 3.5, 0.98, 0.95, [(name[:4], 9, MUTED, True), (name[4:], 8, MUTED, False)], LINE)
    ln = s.shapes.add_connector(2, E(1.34 + i * 1.05), E(3.5), E(2.4), E(3.0))
    ln.line.color.rgb = LINE; ln.line.width = Pt(1)
# SPAN mirror line
a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, E(3.95), E(2.35), E(1.0), E(0.3))
a.fill.solid(); a.fill.fore_color.rgb = GREEN; a.line.fill.background(); a.shadow.inherit = False
txt(s, 3.9, 1.95, 1.4, 0.35, [[("SPAN/TAP", 10, GREEN, True)]], align=PP_ALIGN.CENTER)
# RelayGuard
rect(s, 5.05, 1.65, 3.5, 3.15, fill=CARD, line=GREEN, lw=1.25, round_=True)
txt(s, 5.25, 1.75, 3.1, 0.35, [[("RelayGuard 邊緣閘道", 13, GREEN, True)]])
chip(s, 5.3, 2.2, 3.0, 0.55, [("Mirror NIC ·  唯讀 promiscuous", 10, GREEN, True)], GREEN)
chip(s, 5.3, 2.85, 3.0, 0.55, [("Management NIC ·  管理 / 上報", 10, BLUE, True)], BLUE)
txt(s, 5.3, 3.55, 3.1, 1.1,
    [[("• Mirror NIC 接 SPAN 鏡像埠，純接收", 10.5, MUTED, False)],
     [("• Management NIC 接管理網段", 10.5, MUTED, False)],
     [("• 兩網實體隔離，不對 OT 注入流量", 10.5, MUTED, False)]], sa=2)
# uplink
a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, E(8.6), E(2.9), E(0.6), E(0.28))
a.fill.solid(); a.fill.fore_color.rgb = BLUE; a.line.fill.background(); a.shadow.inherit = False
chip(s, 9.3, 2.4, 3.1, 1.3, [("Control Plane /", 13, WHITE, True), ("SenseL Portal", 13, WHITE, True),
                              ("(管理網段 / 雲端)", 10, MUTED, False)], BLUE)
banner(s, 5.1, ("關鍵：Mirror 埠為唯讀單向鏡像，RelayGuard 不會、也無法對電驛送出封包 — 物理層保證零干擾。", 13, WHITE, True))
txt(s, 0.85, 5.95, 11.6, 0.8,
    [[("交換器設定：", 12.5, GREEN, True),
      ("於 OT 交換器設定 SPAN/Port Mirror，將電驛所在 VLAN/埠的雙向流量鏡像至 Mirror NIC；若為 TAP，使用被動分光/銅纜 TAP。", 12.5, MUTED, False)]])
footer(s, 7)

# ═══════════ Slide 8 — 硬體與規格 ═══════════
s = slide()
header(s, "硬體與規格", "部署需求與效能目標")
cards = [
    ("邊緣硬體", "Lab/PoC：Raspberry Pi 4 (8GB) + SSD；現場高流量建議工業級閘道或 x86。雙網路介面 (內建 + USB Ethernet)。", GREEN),
    ("作業系統", "Raspberry Pi OS 64-bit 或 Ubuntu Server 22.04/24.04；Docker Compose 部署。", BLUE),
    ("網路介面", "Management NIC (管理/上報) + Mirror NIC (SPAN/TAP 唯讀擷取，promiscuous)。", PURPLE),
]
x = 0.85
for t, b, c in cards:
    card(s, x, 1.85, 3.75, 2.0, t, b, c)
    x += 3.9
# perf targets
txt(s, 0.85, 4.15, 11.6, 0.4, [[("效能目標 (MVP)", 15, GREEN, True)]])
perf = [("≤ 100 Mbps", "穩定監控鏡像流量"), ("60 秒", "特徵摘要週期"),
        ("< 3 秒", "規則事件延遲"), ("72 小時", "Lab 連續穩定運行")]
x = 0.85
for v, d in perf:
    chip(s, x, 4.6, 2.78, 1.0, [(v, 18, GREEN, True), (d, 11, MUTED, False)], GREEN)
    x += 2.93
banner(s, 5.95, ("Pi4 適用 Lab / PoC；正式現場依鏡像流量與保留需求選用工業級閘道，並以 SSD 留存 PCAP。", 13, WHITE, True))
footer(s, 8)

# ═══════════ Slide 9 — 配置 (1) 設備納管 ═══════════
s = slide()
header(s, "配置方式 (1)", "電驛資產納管 — Edge Console 設定精靈")
txt(s, 0.85, 1.78, 11.6, 0.5,
    [[("透過 Edge Console (", 13, MUTED, False), ("http://<gw-ip>:8090", 13, GREEN, True),
      (") 完成註冊與設備納管；亦可被動發現鏡像流量中的資產。", 13, MUTED, False)]])
card(s, 0.85, 2.45, 5.65, 1.5, "① 感測器註冊", "設定精靈填入 SenseL API URL、API Key 與企業邀請碼 (registration_token)，綁定 tenant。", GREEN, num=None)
card(s, 6.8, 2.45, 5.65, 1.5, "② 被動資產發現", "從 Mirror 流量自動列出電驛 / IED 的 MAC/IP 與通訊對，與 EdgeX 設備做映射。", BLUE, num=None)
# yaml snippet
rect(s, 0.85, 4.15, 11.6, 2.4, fill=BG2, line=LINE, lw=0.75, round_=True)
txt(s, 1.1, 4.25, 11.0, 0.35, [[("③ EdgeX 設備設定範例 (config/edgex/devices/modbus-relay.yaml)", 13, GREEN, True)]])
code = ("deviceList:\n"
        "  - name: relay-01\n"
        "    profileName: ModbusRelay\n"
        "    protocols:\n"
        "      modbus-tcp:\n"
        "        Address: 192.168.10.20   # 電驛 IP (現場)\n"
        "        Port: \"502\"\n"
        "        UnitID: \"1\"\n"
        "    autoEvents:\n"
        "      - interval: 10s\n"
        "        sourceName: Status        # 主動輪詢狀態 (唯讀)")
txt(s, 1.15, 4.65, 11.0, 1.85, [[(code, 11, WHITE, False)]], sa=0)
footer(s, 9)

print("relay slides 5-9 ready")
