"""Generate the merged Avocado.ai / SenseL ~20-page investor pitch deck.

Merges the existing investor pitch (v5) narrative & design system with the
OT Trust Layer briefing (device mgmt, edge defense, Layer A/B/C, trust engine,
ZTNA/PUF/PQC partner story).

Run: python3 docs/build_pitch_deck.py
Output: docs/Avocado_ai_SenseL_Pitch_v6_OT_TrustLayer.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
ASSET = os.path.join(HERE, "assets")
PITCH = os.path.join(ASSET, "pitch")
LOGO = os.path.join(ASSET, "sensel-edgex-logo.png")
IMG_ARCH = os.path.join(PITCH, "slide5_img4.png")
IMG_OT = os.path.join(PITCH, "slide9_img4.png")
IMG_LOOP = os.path.join(PITCH, "slide11_img4.png")
OUT = os.path.join(HERE, "Avocado_ai_SenseL_Pitch_v6_OT_TrustLayer.pptx")

EMU = 914400
SW, SH = int(13.333 * EMU), int(7.5 * EMU)

# ---- design system (from investor deck v5) ----
NAVY = RGBColor(0x07, 0x16, 0x2D)
LIGHT = RGBColor(0xF6, 0xFA, 0xF8)
GREEN = RGBColor(0x28, 0xA7, 0x45)
BLUE = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
BRAND = RGBColor(0x7D, 0xBA, 0x3B)
SUB = RGBColor(0x64, 0x74, 0x8B)
BODY = RGBColor(0x33, 0x41, 0x55)
BORDER = RGBColor(0xD8, 0xE5, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DSUB = RGBColor(0xC7, 0xD7, 0xE7)
DFOOT = RGBColor(0xA8, 0xBE, 0xD1)
CARDBG = RGBColor(0x11, 0x2B, 0x47)

HEAD = "Aptos Display"
FONT = "Aptos"
ACCENTS = [GREEN, BLUE, PURPLE, AMBER, CYAN]

prs = Presentation()
prs.slide_width = Emu(SW)
prs.slide_height = Emu(SH)
BLANK = prs.slide_layouts[6]


def E(v):
    return Emu(int(v * EMU))


def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(SW), Emu(SH))
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    # top green bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(SW), E(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background(); bar.shadow.inherit = False
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, round_=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(st, E(x), E(y), E(w), E(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        sa=3, font=FONT):
    tb = s.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa); p.space_before = Pt(0)
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = font
    return tb


def fit_pic(s, path, x, y, maxw, maxh, src=(1672, 941)):
    aw, ah = src
    ar = aw / ah
    w = maxw; h = w / ar
    if h > maxh:
        h = maxh; w = h * ar
    px = x + (maxw - w) / 2
    py = y + (maxh - h) / 2
    s.shapes.add_picture(path, E(px), E(py), E(w), E(h))


def header(s, title, subtitle, dark=False):
    tcol = WHITE if dark else NAVY
    scol = DSUB if dark else SUB
    txt(s, 0.55, 0.42, 12.2, 0.5, [[(title, 34, tcol, True)]], font=HEAD)
    txt(s, 0.58, 0.96, 12.0, 0.35, [[(subtitle, 17, scol, False)]])


def footer(s, idx, dark=False):
    fcol = DFOOT if dark else SUB
    txt(s, 0.55, 7.10, 6.0, 0.22,
        [[("Avocado.ai  |  SenseL Investor Pitch", 8.5, fcol, False)]])
    txt(s, 12.2, 7.10, 0.6, 0.22, [[(f"{idx:02d}", 9, fcol, False)]],
        align=PP_ALIGN.RIGHT)


def banner(s, y, text_runs):
    rect(s, 1.05, y, 11.25, 0.62, fill=NAVY, round_=False)
    txt(s, 1.35, y, 10.65, 0.62, [[text_runs]], align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE)


def badge_card(s, x, y, w, h, num, title, body, accent):
    rect(s, x, y, w, h, fill=WHITE, line=BORDER, lw=1.0, round_=True)
    b = rect(s, x + 0.18, y + 0.2, 0.46, 0.46, fill=accent, round_=False)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num)
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = HEAD
    txt(s, x + 0.78, y + 0.2, w - 0.95, 0.4, [[(title, 15, accent, True)]])
    txt(s, x + 0.2, y + 0.78, w - 0.4, h - 0.95, [[(body, 12, BODY, False)]])


# ════════════════════ Slide 1 — Title (navy) ════════════════════
s = slide(bg=NAVY)
rect(s, 0, 7.42, 13.333, 0.08, fill=GREEN)
# right side panel + logo
rect(s, 8.05, 1.0, 4.75, 5.3, fill=CARDBG, line=GREEN, lw=1.0, round_=True)
lw_ = 4.0; lh_ = lw_ * 682 / 1024
s.shapes.add_picture(LOGO, E(8.05 + (4.75 - lw_) / 2), E(2.55), E(lw_), E(lh_))
txt(s, 8.3, 4.55, 4.25, 0.4, [[("AI-native · IT + OT · Trust Layer", 12, DSUB, False)]],
    align=PP_ALIGN.CENTER)
# left text
txt(s, 0.74, 0.74, 4.0, 0.4, [[("Avocado.ai", 22, BRAND, True)]], font=HEAD)
txt(s, 0.72, 1.7, 7.0, 1.5,
    [[("AI-native Security Validation", 38, WHITE, True)],
     [("& OT Trust Platform", 38, WHITE, True)]], font=HEAD, sa=2)
txt(s, 0.76, 3.35, 6.9, 0.9,
    [[("We turn threat intelligence, alerts and AI red-team automation into "
       "verified defense outcomes — across IT and OT.", 17, DSUB, False)]])
rect(s, 0.78, 4.62, 2.8, 0.4, fill=CARDBG, line=GREEN, lw=1.0, round_=True)
txt(s, 0.78, 4.62, 2.8, 0.4, [[("Strategic Angel / Seed Pitch", 11, GREEN, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, 3.7, 4.62, 0.95, 0.4, fill=CARDBG, line=CYAN, lw=1.0, round_=True)
txt(s, 3.7, 4.62, 0.95, 0.4, [[("2026", 11, CYAN, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.78, 5.5, 7.0, 0.4,
    [[("Bigger story:  Detect → Validate → Prove → Improve", 18, WHITE, True)]])
txt(s, 0.55, 7.12, 6.0, 0.22,
    [[("Avocado.ai  |  SenseL Investor Pitch", 8.5, DFOOT, False)]])

# ════════════════════ Slide 2 — Problem ════════════════════
s = slide()
header(s, "Problem", "Security teams are overwhelmed — and OT is the blind spot")
txt(s, 0.7, 1.42, 11.8, 0.5,
    [[("Enterprises don't need another dashboard. They need a validation loop that "
       "connects intelligence, alerts, evidence and improvement — including OT.",
       20, NAVY, True)]])
cards = [
    ("Fragmented signals", "EDR, NDR, WAF, firewall and OT logs do not naturally explain each other.", GREEN),
    ("Non-actionable CTI", "IoCs, CVEs and supply-chain news rarely become testable defense actions.", BLUE),
    ("Unproven controls", "Customers cannot easily prove whether controls detect, block and respond.", PURPLE),
    ("OT trust gap", "OT devices have no hardware identity, no quantum-safe channels, no zero-trust enforcement.", AMBER),
]
x = 0.72
for i, (t, b, c) in enumerate(cards, 1):
    badge_card(s, x, 2.35, 2.82, 2.15, i, t, b, c)
    x += 2.97
banner(s, 5.28, ("Core gap:  \u201cWhat happened?\u201d is not enough. Customers need to know "
                 "\u201cDid our defense actually work?\u201d", 17, WHITE, True))
footer(s, 2)

# ════════════════════ Slide 3 — Why Now ════════════════════
s = slide()
header(s, "Why Now", "AI, validation, and zero-trust OT are converging")
cards = [
    ("AI changes attack speed", "AI accelerates phishing, recon, malware variation and attack automation — defenders need AI agents for analysis and validation.", "AI-assisted security operations", GREEN),
    ("Buyers demand proof", "Budgets rise, but CISOs and boards ask whether EDR, NDR, WAF and OT controls actually work.", "Continuous security validation", BLUE),
    ("OT + Zero-Trust + PQC", "Factories and critical infrastructure face zero-trust mandates and post-quantum migration — with no enterprise SOC to lean on.", "Future-proof OT trust layer", PURPLE),
]
x = 0.72
for i, (t, b, tag, c) in enumerate(cards, 1):
    badge_card(s, x, 1.7, 3.9, 2.9, i, t, b, c)
    rect(s, x + 0.2, 4.15, 3.5, 0.3, fill=None)
    txt(s, x + 0.2, 4.12, 3.5, 0.35, [[(tag, 12, c, True)]])
    x += 4.07
banner(s, 5.35, ("Security is shifting from manual monitoring to AI-driven validation. "
                 "SenseL is built for that shift — IT and OT.", 17, WHITE, True))
footer(s, 3)

# ════════════════════ Slide 4 — Solution ════════════════════
s = slide()
header(s, "Solution", "SenseL: AI-native Security Validation, extended to OT")
txt(s, 0.7, 1.42, 11.8, 0.5,
    [[("SenseL connects threat intelligence, telemetry and AI red-team automation to "
       "continuously validate whether enterprise and OT defenses actually work.",
       20, NAVY, True)]])
steps = [
    ("Ingest", "EDR / NDR / WAF / OT telemetry / firewall / CTI / CVE / IoC / news", GREEN),
    ("Reason", "AI agents correlate alerts, generate hypotheses and map scenarios to MITRE ATT&CK", BLUE),
    ("Validate", "AI Red Team + BAS workflows test whether controls detect, block and respond", PURPLE),
    ("Deliver", "Evidence-linked reports, defense gaps, remediation playbooks and executive summaries", AMBER),
]
x = 0.72
for i, (t, b, c) in enumerate(steps, 1):
    badge_card(s, x, 2.35, 2.82, 2.15, i, t, b, c)
    if i < 4:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, E(x + 2.84), E(3.25), E(0.12), E(0.22))
        a.fill.solid(); a.fill.fore_color.rgb = SUB; a.line.fill.background(); a.shadow.inherit = False
    x += 2.97
banner(s, 5.28, ("SenseL helps customers move from security monitoring to verified defense "
                 "outcomes — now including OT trust.", 17, WHITE, True))
footer(s, 4)

# ════════════════════ Slide 5 — Product Architecture (image) ════════════════════
s = slide()
header(s, "Product Architecture", "Telemetry → Intelligence → Validation → Operations")
fit_pic(s, IMG_ARCH, 0.55, 1.3, 12.25, 5.0)
banner(s, 6.42, ("One platform connects existing tools, CTI, AI red-team automation and "
                 "customer-ready operations.", 14.5, WHITE, True))
footer(s, 5)

# ════════════════════ Slide 6 — Commercial Wedge ════════════════════
s = slide()
header(s, "Commercial Wedge", "Land with WAF, expand with XDR, scale into OT validation")
cards = [
    ("WAF MSSP", "Fastest revenue entry. Radware WAF integration + DNS onboarding for SMBs, schools and websites.", "Low-friction onboarding · Recurring revenue", GREEN),
    ("XDR / AI SOC", "Core platform expansion. EDR, NDR, firewall, CTI and AI reporting for managed security operations.", "Monthly managed service · Exec reporting", BLUE),
    ("OT Validation", "Premium vertical. Energy & equipment telemetry becomes OT evidence, anomaly validation and trust assurance.", "Higher-value accounts · Strong differentiation", PURPLE),
]
x = 0.72
for i, (t, b, tag, c) in enumerate(cards, 1):
    badge_card(s, x, 1.7, 3.9, 2.9, i, t, b, c)
    txt(s, x + 0.2, 4.12, 3.5, 0.5, [[(tag, 11.5, c, True)]])
    x += 4.07
banner(s, 5.35, ("WAF gets us in. XDR expands the account. OT creates premium "
                 "differentiation and a defensible trust moat.", 17, WHITE, True))
footer(s, 6)

# ════════════════════ Slide 7 — SenseL OT Layer A/B/C (image) ════════════════════
s = slide()
header(s, "SenseL OT — Layer A / B / C", "From industrial data → AI insights → validated, trusted actions")
fit_pic(s, IMG_OT, 0.55, 1.3, 12.25, 5.0)
banner(s, 6.42, ("OT thesis: start with measurable energy / ESG value, then expand into "
                 "evidence-based OT resilience and security validation.", 14.5, WHITE, True))
footer(s, 7)

# ════════════════════ Slide 8 — OT Edge & Device Trust (NEW) ════════════════════
s = slide()
header(s, "OT Edge & Device Trust", "Passive-first edge sensor — see everything, disrupt nothing")
lw_ = 1.7; lh_ = lw_ * 682 / 1024
s.shapes.add_picture(LOGO, E(12.85 - lw_ - 0.05), E(0.42), E(lw_), E(lh_))
cards = [
    ("Device Management", "Three discovery tracks: sensor self-registration, EdgeX active devices (Modbus/OPC-UA/S7), and zero-risk passive mirror discovery.", GREEN),
    ("Device-Side Defense", "OT-001~019 detection: network behavior, IEC 61850 (GOOSE/MMS) and CTI IoC matching — with evidence capture.", BLUE),
    ("CTI & Rule Distribution", "Cloud produces signed policy; edge subscribes over MQTT/HTTP and hot-reloads blacklists, rules and baselines.", PURPLE),
]
x = 0.72
for i, (t, b, c) in enumerate(cards, 1):
    badge_card(s, x, 1.85, 3.9, 2.55, i, t, b, c)
    x += 4.07
banner(s, 4.95, ("Deployed on industrial gateways / Raspberry Pi. Passive by design — "
                 "no process disruption, full visibility.", 16, WHITE, True))
txt(s, 0.72, 5.85, 11.8, 0.9,
    [[("Edge runs offline-resilient with local buffering. Today the edge alerts and "
       "recommends; enforcement is the partner-enabled next step (ZTNA).", 13, BODY, False)]])
footer(s, 8)

# ════════════════════ Slide 9 — The Trust Engine (NEW) ════════════════════
s = slide()
header(s, "The Trust Engine", "Explainable, behavior-based trust for every OT entity")
txt(s, 0.7, 1.42, 11.8, 0.5,
    [[("Layer B computes a continuous, auditable trust score per entity — the core of "
       "our OT differentiation.", 18, NAVY, True)]])
# trust level chips
levels = [("Healthy", "\u2265 0.85", GREEN), ("Suspicious", "\u2265 0.70", CYAN),
          ("Degraded", "\u2265 0.50", AMBER), ("Critical", "< 0.50", PURPLE)]
x = 0.72
for name, rng, c in levels:
    rect(s, x, 2.15, 2.82, 0.8, fill=WHITE, line=c, lw=1.5, round_=True)
    txt(s, x, 2.22, 2.82, 0.35, [[(name, 15, c, True)]], align=PP_ALIGN.CENTER)
    txt(s, x, 2.58, 2.82, 0.3, [[(rng, 12, BODY, False)]], align=PP_ALIGN.CENTER)
    x += 2.97
# mechanism cards
mechs = [
    ("Penalty / Recovery", "Trust drops on malicious episodes (severity · recurrence · sequence anomaly · confidence) and recovers slowly when clean.", GREEN),
    ("Risk Propagation", "A degraded/critical entity spreads risk along a trust graph (identity → workstation → PLC), fully audited.", BLUE),
    ("Trust Alerts", "Discrete alerts on critical drops, multi high-risk labels, and C2 + exfiltration combinations.", PURPLE),
]
x = 0.72
for i, (t, b, c) in enumerate(mechs, 1):
    badge_card(s, x, 3.2, 3.9, 2.0, i, t, b, c)
    x += 4.07
banner(s, 5.5, ("Not \u201cLLM over logs.\u201d A formula-driven, explainable trust engine "
                "that security and OT teams can audit.", 16, WHITE, True))
footer(s, 9)

# ════════════════════ Slide 10 — OT Trust Layer & Partners (NEW) ════════════════════
s = slide()
header(s, "OT Trust Layer & Strategic Partners", "Identity + Channel + Access + Behavior = one Unified Trust Score")
cols = [
    ("PUF", "Identity", "Hardware root of identity — unclonable device ID, key wrapping, attestation.", PURPLE),
    ("PQC", "Channel", "Quantum-safe channels & signatures — mTLS and end-to-end signed policy.", BLUE),
    ("ZTNA", "Access", "Device posture, dynamic authorization, micro-segmentation and enforcement.", AMBER),
    ("SenseL", "Behavior", "Existing behavioral trust engine — entity state, alerts, propagation.", GREEN),
]
x = 0.72
for label, role, body, c in cols:
    rect(s, x, 1.7, 2.82, 3.0, fill=WHITE, line=c, lw=1.5, round_=True)
    txt(s, x, 1.9, 2.82, 0.45, [[(label, 20, c, True)]], align=PP_ALIGN.CENTER, font=HEAD)
    txt(s, x, 2.4, 2.82, 0.3, [[(role, 13, SUB, True)]], align=PP_ALIGN.CENTER)
    txt(s, x + 0.2, 2.85, 2.42, 1.7, [[(body, 12, BODY, False)]], align=PP_ALIGN.CENTER)
    x += 2.97
# unified score bar
rect(s, 1.05, 5.05, 11.25, 0.85, fill=NAVY, round_=True)
txt(s, 1.35, 5.05, 10.65, 0.85,
    [[("Fused into one ", 16, WHITE, True), ("Unified Trust Score", 17, BRAND, True),
      (" that drives real-time access and isolation decisions.", 16, WHITE, True)]],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.72, 6.15, 11.8, 0.8,
    [[("Partner play:  ", 13, GREEN, True),
      ("PUF, PQC and ZTNA partners migrate into the SenseL Trust Layer — making OT "
       "validation defensible, future-proof and standards-aligned (NIST PQC).", 13, BODY, False)]])
footer(s, 10)

# ════════════════════ Slide 11 — AI Red Team Validation Loop (image) ════════════════════
s = slide()
header(s, "AI Red Team Validation Loop", "From threat intelligence to verified defense")
fit_pic(s, IMG_LOOP, 1.8, 1.35, 9.7, 4.85)
banner(s, 6.42, ("Every new threat becomes a testable scenario, and every test creates "
                 "evidence — the engine that powers SenseL.", 14.5, WHITE, True))
footer(s, 11)

# ════════════════════ Slide 12 — Technical Moat ════════════════════
s = slide()
header(s, "Technical Moat", "Evidence Chain + Agentic Eval + OT Trust Layer")
txt(s, 0.7, 1.42, 11.8, 0.5,
    [[("The moat is the evidence, evaluation and trust loop around the model — not the model itself.",
       18, NAVY, True)]])
cards = [
    ("Evidence Chain", "Every alert, IoC, hypothesis, scenario, telemetry record and action is linked into traceable evidence — prove what happened and what was tested.", GREEN),
    ("Agentic Eval", "AI agents, red-team workflows and defense outcomes are evaluated for accuracy, safety and reliability — reducing hallucination and unsafe automation.", BLUE),
    ("OT Trust Layer", "Explainable trust scoring + PUF/PQC/ZTNA integration creates a defensible, standards-aligned OT moat competitors can't copy quickly.", PURPLE),
]
x = 0.72
for i, (t, b, c) in enumerate(cards, 1):
    badge_card(s, x, 2.15, 3.9, 2.75, i, t, b, c)
    x += 4.07
banner(s, 5.45, ("Our moat is the evidence and evaluation loop — and an OT trust layer "
                 "that gets stronger with every deployment.", 16, WHITE, True))
footer(s, 12)

# ════════════════════ Slide 13 — Traction: WAF MSSP ════════════════════
s = slide()
header(s, "Traction — WAF MSSP", "Revenue-ready wedge with Radware / DNS onboarding")
txt(s, 0.7, 1.42, 11.8, 0.45,
    [[("WAF is the fastest near-term entry: clear pain, light onboarding, recurring service.",
       16, NAVY, True)]])
items = [("Radware", "WAF API integration path", GREEN),
         ("DNS", "CNAME / A record onboarding", BLUE),
         ("MVP", "Multi-tenant UI + mock API", PURPLE),
         ("3–10", "Initial website targets", AMBER)]
x = 0.72
for t, b, c in items:
    rect(s, x, 2.05, 2.82, 1.15, fill=WHITE, line=BORDER, lw=1.0, round_=True)
    txt(s, x, 2.2, 2.82, 0.5, [[(t, 20, c, True)]], align=PP_ALIGN.CENTER, font=HEAD)
    txt(s, x + 0.1, 2.72, 2.62, 0.4, [[(b, 11.5, BODY, False)]], align=PP_ALIGN.CENTER)
    x += 2.97
badge_card(s, 0.72, 3.45, 5.85, 1.7, 1, "Current progress",
           "Radware WAF API / Postman reference reviewed; DNS onboarding and multi-tenant "
           "service workflow in implementation for MSSP operation.", GREEN)
badge_card(s, 6.77, 3.45, 5.85, 1.7, 2, "Evidence to add",
           "Partner letter, first customer list, demo screenshots, onboarding time, WAF "
           "event examples, pricing sheet, conversion timeline.", BLUE)
banner(s, 5.45, ("Investor takeaway: WAF creates the first repeatable commercial package "
                 "and validates partner-led go-to-market.", 16, WHITE, True))
footer(s, 13)

# ════════════════════ Slide 14 — Traction: XDR / AI SOC ════════════════════
s = slide()
header(s, "Traction — XDR / AI SOC", "Core platform validation across EDR, NDR and AI reporting")
rows = [
    ("Tung's Hospital", "Wazuh EDR", "June 2026 target: >500 endpoints", "Healthcare XDR reference", GREEN),
    ("\u9f0e\u65b0 \u667a\u7ba1\u5bb6", "SentinelOne + FortiSIEM", "XDR with enterprise software channel", "Retail / ERP distribution", BLUE),
    ("Qsec / Arista", "NDR integration", "NDR integrates with Arista telemetry", "Network security validation", PURPLE),
    ("SenseL Reports", "AI executive summary", "EDR / NDR / CTI evidence into monthly report", "AI SOC automation", AMBER),
]
# header row
hx = [0.72, 3.4, 6.0, 9.3]; hw = [2.6, 2.5, 3.2, 3.0]
heads = ["Account / Partner", "Technology", "Status / Evidence", "Commercial signal"]
for i, h in enumerate(heads):
    rect(s, hx[i], 1.65, hw[i], 0.45, fill=NAVY)
    txt(s, hx[i] + 0.1, 1.65, hw[i] - 0.2, 0.45, [[(h, 11.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
y = 2.1
for acct, tech, st, sig, c in rows:
    rect(s, 0.72, y, 11.58, 0.86, fill=WHITE, line=BORDER, lw=0.75)
    rect(s, 0.72, y, 0.1, 0.86, fill=c)
    txt(s, 0.92, y, 2.45, 0.86, [[(acct, 13, c, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 3.5, y, 2.4, 0.86, [[(tech, 12, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 6.1, y, 3.1, 0.86, [[(st, 11.5, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 9.4, y, 2.85, 0.86, [[(sig, 11.5, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.94
banner(s, y + 0.05, ("XDR / AI SOC turns field deployments into repeatable managed operations, "
                     "feeding EDR/NDR/CTI evidence into SenseL reports.", 14.5, WHITE, True))
footer(s, 14)

# ════════════════════ Slide 15 — Traction: SenseL OT ════════════════════
s = slide()
header(s, "Traction — SenseL OT", "Energy data is the entry; OT security validation is the expansion")
cards = [
    ("Smart Meter PoC", "Entered a semiconductor-line smart clamp meter ESG + OT analytics PoC. Q3 completion planned.", GREEN),
    ("OT Security Expansion", "Q4 expansion from energy/ESG analytics toward OT cybersecurity validation and trust scoring.", BLUE),
    ("Edge + Trust Layer", "RelayGuard edge sensor + Layer A/B/C trust engine + PUF/PQC/ZTNA partner roadmap.", PURPLE),
]
x = 0.72
for i, (t, b, c) in enumerate(cards, 1):
    badge_card(s, x, 1.85, 3.9, 2.55, i, t, b, c)
    x += 4.07
banner(s, 4.95, ("OT vertical thesis: measurable energy / ESG value first, then evidence-based "
                 "OT resilience and security validation.", 16, WHITE, True))
txt(s, 0.72, 5.85, 11.8, 0.9,
    [[("Evidence to add:  ", 13, GREEN, True),
      ("PoC scope, smart clamp meter architecture, Q3 acceptance criteria, Q4 OT security "
       "expansion plan, and partner LOIs for the trust layer.", 13, BODY, False)]])
footer(s, 15)

# ════════════════════ Slide 16 — Traction: 2026 Momentum ════════════════════
s = slide()
header(s, "Traction — 2026 Momentum", "Concrete milestones turning vision into Taiwan market entry")
cards = [
    ("Hsinchu Industry SIG", "Cooperate with Hsinchu Industrial Union on a cybersecurity SIG plan (Q3/Q4 2026); SenseL for supply-chain SMB adoption.", "Go-to-market access", GREEN),
    ("MoDA Cybersecurity Entry", "Applying for cybersecurity capability registration and self-developed R&D recognition for Taiwan market entry.", "Market qualification", BLUE),
    ("OT Smart Meter PoC", "Semiconductor-line smart clamp meter ESG + OT analytics PoC; Q3 completion, Q4 OT cybersecurity expansion.", "Premium vertical proof", PURPLE),
    ("MOEA SBIR Submission", "R&D proposal submitted to MOEA SBIR, targeting Q3 execution to support the SenseL validation engine.", "Non-dilutive leverage", AMBER),
]
x = 0.72
for i, (t, b, tag, c) in enumerate(cards, 1):
    badge_card(s, x, 1.75, 2.82, 3.0, i, t, b, c)
    txt(s, x + 0.2, 4.4, 2.5, 0.3, [[(tag, 11, c, True)]])
    x += 2.97
banner(s, 5.45, ("Signal: early access to supply-chain SMBs, OT semiconductor use cases, "
                 "Taiwan cybersecurity programs and SBIR R&D support.", 15, WHITE, True))
footer(s, 16)

# ════════════════════ Slide 17 — Business Model ════════════════════
s = slide()
header(s, "Business Model", "One platform, multiple recurring revenue streams")
cards = [
    ("Platform Subscription", "Tenants, protected assets, endpoints, websites and OT data sources.", "Monthly / annual SaaS", GREEN),
    ("Managed Security Service", "AI-assisted WAF MSSP, XDR monitoring, incident summary and monthly report.", "Recurring service fee", BLUE),
    ("Validation-as-a-Service", "AI Red Team validation, BAS scenarios, defense effectiveness and evidence reports.", "Per assessment / package", PURPLE),
    ("Partner Revenue Share", "SI, MSSP, WAF providers, OT and trust-layer partners (PUF/PQC/ZTNA).", "Revenue share / reseller margin", AMBER),
]
x = 0.72
for i, (t, b, tag, c) in enumerate(cards, 1):
    badge_card(s, x, 1.75, 2.82, 2.85, i, t, b, c)
    txt(s, x + 0.2, 4.25, 2.5, 0.3, [[(tag, 11, c, True)]])
    x += 2.97
banner(s, 5.35, ("Land through services, expand through subscriptions, scale through partners.",
                 17, WHITE, True))
footer(s, 17)

# ════════════════════ Slide 18 — Go-to-Market ════════════════════
s = slide()
header(s, "Go-to-Market", "Partner-led, vertical-focused, validation-driven")
rows = [
    ("Q3 2026", "SMB / school / website owners", "WAF MSSP + SIG SMB package", "Fast recurring revenue + SMB access", GREEN),
    ("Q3–Q4 2026", "Healthcare / retail / SMB", "XDR + AI SOC + Red Team validation", "Managed security + closed-loop validation", BLUE),
    ("Q3–Q4 2026", "Factory / semiconductor OT", "SenseL OT smart meter PoC", "Complete OT analytics, expand to OT security", PURPLE),
    ("12–18 mo", "SI / MSSP / industry partners", "Multi-tenant SenseL + Trust Layer", "Channel scale and partner-led delivery", AMBER),
]
heads = ["When", "Who", "What", "Why it matters"]
hx = [0.72, 2.4, 5.6, 9.0]; hw = [1.6, 3.1, 3.3, 3.3]
for i, h in enumerate(heads):
    rect(s, hx[i], 1.6, hw[i], 0.45, fill=NAVY)
    txt(s, hx[i] + 0.1, 1.6, hw[i] - 0.2, 0.45, [[(h, 11.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
y = 2.05
for when, who, what, why, c in rows:
    rect(s, 0.72, y, 11.58, 0.82, fill=WHITE, line=BORDER, lw=0.75)
    rect(s, 0.72, y, 0.1, 0.82, fill=c)
    txt(s, 0.92, y, 1.5, 0.82, [[(when, 12, c, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 2.5, y, 3.0, 0.82, [[(who, 12, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 5.7, y, 3.2, 0.82, [[(what, 12, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 9.1, y, 3.15, 0.82, [[(why, 11.5, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.9
banner(s, y + 0.05, ("We land with urgent security needs, add AI Red Team validation, and scale "
                     "through trusted partners.", 15, WHITE, True))
footer(s, 18)

# ════════════════════ Slide 19 — Fundraising Ask ════════════════════
s = slide()
header(s, "Fundraising Ask", "Strategic Angel Bridge → Seed-ready execution")
txt(s, 0.7, 1.42, 11.8, 0.5,
    [[("Raise ", 18, NAVY, True), ("NT$3M–5M", 18, GREEN, True),
      (" as a Strategic Angel Bridge, then ", 18, NAVY, True),
      ("NT$20M–30M Seed", 18, GREEN, True),
      (" after commercial validation.", 18, NAVY, True)]])
cards = [
    ("Bridge Round", "NT$3M–5M. Complete MVP, convert first commercial pilots, produce evidence assets, prepare seed data room.", GREEN),
    ("Use of Funds", "Productization · MSSP onboarding · WAF/XDR/OT connectors · AI validation engine · OT Trust Layer · GTM.", BLUE),
    ("18-month Milestones", "10–20 paying customers/pilots · multi-tenant platform · WAF MSSP · XDR/AI SOC reports · OT validation + trust demo.", PURPLE),
]
x = 0.72
for i, (t, b, c) in enumerate(cards, 1):
    badge_card(s, x, 2.3, 3.9, 2.6, i, t, b, c)
    x += 4.07
banner(s, 5.45, ("Goal: turn field PoCs into repeatable product modules, recurring revenue "
                 "and Seed-ready traction.", 16, WHITE, True))
footer(s, 19)

# ════════════════════ Slide 20 — Closing / Vision (navy) ════════════════════
s = slide(bg=NAVY)
rect(s, 0, 7.42, 13.333, 0.08, fill=GREEN)
lw_ = 4.4; lh_ = lw_ * 682 / 1024
s.shapes.add_picture(LOGO, E((13.333 - lw_) / 2), E(1.5), E(lw_), E(lh_))
txt(s, 1.0, 4.35, 11.333, 0.6,
    [[("Detect → Validate → Prove → Improve", 28, WHITE, True)]],
    align=PP_ALIGN.CENTER, font=HEAD)
txt(s, 1.0, 5.15, 11.333, 0.5,
    [[("Now extended with an OT Trust Layer:  Identity (PUF) · Channel (PQC) · "
       "Access (ZTNA) · Behavior (SenseL)", 14, DSUB, False)]],
    align=PP_ALIGN.CENTER)
txt(s, 1.0, 6.15, 11.333, 0.4,
    [[("Avocado.ai  ·  SenseL EdgeX  ·  Strategic Angel / Seed 2026", 13, BRAND, True)]],
    align=PP_ALIGN.CENTER)

prs.save(OUT)
print("Saved:", OUT, "slides:", len(prs.slides._sldIdLst))
