"""Generate the SenseL OT Trust Layer partner briefing as a branded PPTX.

Run: python3 docs/build_partner_deck.py
Output: docs/SenseL-OT-Trust-Layer-Briefing.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "assets", "sensel-edgex-logo.png")
OUT = os.path.join(HERE, "SenseL-OT-Trust-Layer-Briefing.pptx")

# 16:9 canvas
EMU = 914400
SW, SH = int(13.333 * EMU), int(7.5 * EMU)

# Brand palette (from logo)
BG = RGBColor(0x0E, 0x16, 0x21)        # deep navy
BG2 = RGBColor(0x16, 0x20, 0x2E)       # panel
GREEN = RGBColor(0x9B, 0xD5, 0x34)     # avocado / X green
PURPLE = RGBColor(0xA2, 0x4B, 0xD0)    # octopus purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9C, 0xA8, 0xB8)
LINE = RGBColor(0x2A, 0x36, 0x48)
CARD = RGBColor(0x1A, 0x25, 0x36)

FONT = "Arial"

prs = Presentation()
prs.slide_width = Emu(SW)
prs.slide_height = Emu(SH)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(SW), Emu(SH))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shape_type, Emu(int(x * EMU)), Emu(int(y * EMU)),
                            Emu(int(w * EMU)), Emu(int(h * EMU)))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4):
    tb = s.shapes.add_textbox(Emu(int(x * EMU)), Emu(int(y * EMU)),
                              Emu(int(w * EMU)), Emu(int(h * EMU)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # runs: list of paragraphs; each paragraph is list of (text, size, color, bold)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (t, size, color, bold) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = FONT
    return tb


def logo_corner(s):
    # top-right small logo
    w = 2.05
    h = w * 682 / 1024
    s.shapes.add_picture(LOGO, Emu(int((13.333 - w - 0.35) * EMU)),
                         Emu(int(0.30 * EMU)), Emu(int(w * EMU)), Emu(int(h * EMU)))


def header(s, kicker, title):
    box(s, 0.55, 0.62, 0.12, 0.92, fill=GREEN)
    text(s, 0.85, 0.55, 9.5, 0.45, [[(kicker, 13, GREEN, True)]])
    text(s, 0.83, 0.92, 10.0, 0.8, [[(title, 30, WHITE, True)]])
    logo_corner(s)
    box(s, 0.85, 1.62, 11.6, 0.018, fill=LINE)


def footer(s, idx):
    text(s, 0.55, 7.04, 8.0, 0.35,
         [[("SenseL EdgeX  ·  by AvocadoAI  ·  Confidential", 9, MUTED, False)]])
    text(s, 11.6, 7.04, 1.2, 0.35, [[(str(idx), 10, MUTED, True)]],
         align=PP_ALIGN.RIGHT)


def table(s, x, y, w, headers, rows, col_w, row_h=0.42, head_h=0.46,
          fs=11, head_fs=11):
    """Lightweight table built from shapes for full color control."""
    cy = y
    # header row
    cx = x
    for j, htext in enumerate(headers):
        cell = box(s, cx, cy, col_w[j], head_h, fill=PURPLE)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(6)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = htext
        r.font.size = Pt(head_fs)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT
        cx += col_w[j]
    cy += head_h
    # body rows
    for ri, row in enumerate(rows):
        cx = x
        fill = CARD if ri % 2 == 0 else BG2
        for j, val in enumerate(row):
            cell = box(s, cx, cy, col_w[j], row_h, fill=fill,
                       line=LINE, line_w=0.5)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Pt(6)
            tf.margin_right = Pt(6)
            tf.margin_top = Pt(1)
            tf.margin_bottom = Pt(1)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            # allow (text, color, bold) tuple or plain string
            if isinstance(val, tuple):
                vtext, vcolor, vbold = val
            else:
                vtext, vcolor, vbold = val, WHITE, False
            r = p.add_run()
            r.text = vtext
            r.font.size = Pt(fs)
            r.font.color.rgb = vcolor
            r.font.bold = vbold
            r.font.name = FONT
            cx += col_w[j]
        cy += row_h
    return cy


def chip(s, x, y, w, h, label, color):
    c = box(s, x, y, w, h, fill=BG2, line=color, line_w=1.5, round_=True)
    tf = c.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    for i, para in enumerate(label):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        for (t, size, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = FONT
    return c


def arrow(s, x, y, w, color=GREEN):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(int(x * EMU)),
                           Emu(int(y * EMU)), Emu(int(w * EMU)), Emu(int(0.22 * EMU)))
    a.shadow.inherit = False
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


def down_arrow(s, x, y, h, color=GREEN):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Emu(int(x * EMU)),
                           Emu(int(y * EMU)), Emu(int(0.22 * EMU)), Emu(int(h * EMU)))
    a.shadow.inherit = False
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


# ───────────────────────────────────────── Slide 1 — Title
s = slide()
# subtle top + bottom accent bars
box(s, 0, 0, 13.333, 0.10, fill=GREEN)
box(s, 0, 7.40, 13.333, 0.10, fill=PURPLE)
lw = 6.2
lh = lw * 682 / 1024
s.shapes.add_picture(LOGO, Emu(int((13.333 - lw) / 2 * EMU)),
                     Emu(int(1.15 * EMU)), Emu(int(lw * EMU)), Emu(int(lh * EMU)))
text(s, 1.0, 4.95, 11.333, 0.7,
     [[("OT Security — Trust Layer Briefing", 30, WHITE, True)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 5.62, 11.333, 0.5,
     [[("Partner Briefing for  ", 16, MUTED, False),
       ("ZTNA · PUF · PQC", 16, GREEN, True),
       ("  Collaboration", 16, MUTED, False)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 6.55, 11.333, 0.4,
     [[("Draft v0.1   ·   Confidential", 12, MUTED, False)]],
     align=PP_ALIGN.CENTER)

# ───────────────────────────────────────── Slide 2 — Executive Summary
s = slide()
header(s, "EXECUTIVE SUMMARY", "Four Trust Questions, Four Partners")
text(s, 0.85, 1.78, 11.6, 0.8,
     [[("SenseL runs a mature ", 14, MUTED, False),
       ("behavioral trust layer", 14, GREEN, True),
       (" today. To reach a complete Zero-Trust posture for OT, three partners "
        "migrate their capabilities into one unified Trust Layer.", 14, MUTED, False)]])
table(s, 0.85, 2.7, 11.6,
      ["Partner", "Provides", "Answers the question"],
      [[("PUF", PURPLE, True), "Hardware root of identity", "Who is this device — really?"],
       [("PQC", PURPLE, True), "Quantum-resistant channels & signing", "Can we trust how it was transmitted?"],
       [("ZTNA", PURPLE, True), "Access decisions & enforcement", "What is this entity allowed to do?"],
       [("SenseL (today)", GREEN, True), "Behavioral trust scoring", "Is it behaving correctly?"]],
      col_w=[2.4, 4.4, 4.8], row_h=0.62, head_h=0.5, fs=12)
chip(s, 0.85, 6.05, 11.6, 0.62,
     [[("The ask:  combine these four signals into one Unified Trust Score that drives "
        "real-time access and isolation decisions.", 13, WHITE, True)]], GREEN)
footer(s, 2)

# ───────────────────────────────────────── Slide 3 — Three Building Blocks
s = slide()
header(s, "WHO WE ARE", "Three Building Blocks")
layers = [
    ("CLOUD CONTROL PLANE", "CTI hub · Policy production · Southbound rule distribution · Multi-agent workforce · Customer portal", PURPLE),
    ("INFERENCE PLANE  (Layer A / B / C)", "Normalization · ET-BERT inference · Episode aggregation · >> TRUST ENGINE << · Agentic investigation", GREEN),
    ("EDGE SENSOR  (industrial gateway / Raspberry Pi)", "Passive traffic detection + active telemetry · Local alerting · Evidence capture · Policy sync", WHITE),
]
y = 1.95
for title_t, desc, col in layers:
    box(s, 0.85, y, 0.14, 1.35, fill=col)
    box(s, 1.05, y, 11.4, 1.35, fill=BG2, line=LINE, line_w=0.75, round_=True)
    text(s, 1.35, y + 0.16, 10.9, 0.45, [[(title_t, 16, col if col != WHITE else WHITE, True)]])
    text(s, 1.35, y + 0.66, 10.9, 0.6, [[(desc, 12.5, MUTED, False)]])
    y += 1.55
footer(s, 3)

# ───────────────────────────────────────── Slide 4 — End-to-End Architecture
s = slide()
header(s, "ARCHITECTURE", "End-to-End Architecture (Purdue-Aligned)")
# four stacked zones with flow
zones = [
    ("FIELD ZONE  (Purdue L1/L2)", "PLC · RTU · IED · HMI", MUTED),
    ("EDGE SENSOR", "Passive detection (SPAN/TAP)  +  Active telemetry (EdgeX)  →  Edge Agent (register · health · policy)", WHITE),
    ("INFERENCE PLANE  (L3)", "Layer A Normalization  →  Layer B Inference + Trust Engine  →  Layer C Agentic Investigation", GREEN),
    ("CLOUD CONTROL PLANE  (L4/L5)", "CTI Hub  →  Policy Engine  →  Southbound Distribution (signed policy + IoC)", PURPLE),
]
y = 1.85
for i, (t, d, col) in enumerate(zones):
    box(s, 1.6, y, 10.1, 0.92, fill=BG2, line=col, line_w=1.25, round_=True)
    text(s, 1.85, y + 0.10, 9.6, 0.4, [[(t, 13, col, True)]])
    text(s, 1.85, y + 0.46, 9.6, 0.4, [[(d, 11, MUTED, False)]])
    if i < 3:
        down_arrow(s, 6.55, y + 0.93, 0.28, color=GREEN if i % 2 == 0 else PURPLE)
    y += 1.22
text(s, 0.85, 6.78, 11.6, 0.4,
     [[("Design principles:  passive-first (no process disruption) · explainable trust · clean north/south separation",
        11, GREEN, True)]], align=PP_ALIGN.CENTER)
footer(s, 4)

# ───────────────────────────────────────── Slide 5 — Device Management
s = slide()
header(s, "CAPABILITY 1", "Device Management")
table(s, 0.85, 1.85, 11.6,
      ["Track", "How", "Process risk"],
      [["Sensor registration", "Edge agent registers to cloud, bound to a tenant", "—"],
       ["Active devices", "Guided wizard: Modbus / MQTT / OPC-UA / S7", ("Low (read-only)", GREEN, False)],
       ["Passive discovery", "Mirror traffic reveals MAC/IP, mapped to devices", ("Zero (passive)", GREEN, False)]],
      col_w=[2.9, 6.2, 2.5], row_h=0.56, head_h=0.48, fs=12)
text(s, 0.85, 4.05, 11.6, 0.6,
     [[("Plus optional read-only active probing (never writes to a device) and behavioral baseline learning.",
        12.5, MUTED, False)]])
chip(s, 0.85, 5.6, 11.6, 0.95,
     [[("GAP → PUF", 14, PURPLE, True)],
      [("Device identity relies today on a copyable API key. There is no hardware-bound, unclonable identity.",
        13, WHITE, False)]], PURPLE)
footer(s, 5)

# ───────────────────────────────────────── Slide 6 — Device-Side Defense
s = slide()
header(s, "CAPABILITY 2", "Device-Side Defense")
table(s, 0.85, 1.85, 11.6,
      ["Capability", "What it does"],
      [["Detection engine", "OT network-behavior rules, IEC 61850 (GOOSE/MMS), CTI IoC hits"],
       ["Threat-intel matching", "Pulls/receives blacklists; hot-reloads local IoC cache"],
       ["Policy enforcement", "Applies cloud-pushed enabled rules + baselines"],
       ["Evidence capture", "In-memory ring buffer with evidence references"],
       ["Offline resilience", "Local buffering + exponential-backoff reconnect"]],
      col_w=[3.4, 8.2], row_h=0.5, head_h=0.46, fs=12)
chip(s, 0.85, 5.05, 11.6, 0.62,
     [[("Today:  \u201cBlock / Quarantine\u201d are recommended actions only — the edge does not block inline.",
        12.5, WHITE, True)]], GREEN)
chip(s, 0.85, 5.85, 11.6, 0.78,
     [[("GAP → ZTNA", 14, PURPLE, True)],
      [("No \u201ctrust score → dynamic access control / micro-segmentation\u201d enforcement.", 13, WHITE, False)]],
     PURPLE)
footer(s, 6)

# ───────────────────────────────────────── Slide 7 — Layer A/B/C + Trust Engine
s = slide()
header(s, "CAPABILITY 3", "Layer A / B / C Analytics  ·  The Trust Engine")
# pipeline chips
px = 0.85
labels = [("Layer A\nNormalize", GREEN), ("Layer B-1\nET-BERT Cascade", GREEN),
          ("Layer B-2\nEpisodes", GREEN), ("TRUST ENGINE\nentity score", PURPLE),
          ("Layer C\nInvestigation", GREEN)]
cw = 2.05
for i, (lab, col) in enumerate(labels):
    parts = lab.split("\n")
    chip(s, px, 1.85, cw, 0.85,
         [[(parts[0], 12, col, True)], [(parts[1], 10, MUTED, False)]], col)
    if i < len(labels) - 1:
        arrow(s, px + cw + 0.02, 2.18, 0.18, color=col)
    px += cw + 0.22
# trust engine detail card
box(s, 0.85, 3.05, 11.6, 3.4, fill=BG2, line=LINE, line_w=0.75, round_=True)
text(s, 1.15, 3.2, 11.0, 0.45,
     [[("The Trust Engine — our existing software trust layer", 16, GREEN, True)]])
bullets = [
    ("Formula:  ", "T_new = clamp( T_old − Penalty(event) + Recovery(time) )"),
    ("Penalty:  ", "weighs severity · recurrence · sequence anomaly · classifier confidence"),
    ("Recovery:  ", "slow, only when no new malicious activity occurs"),
    ("Levels:  ", "Healthy → Suspicious → Degraded → Critical"),
    ("Propagation:  ", "degraded/critical entity spreads risk along a trust graph (identity → workstation → PLC), fully audited"),
]
by = 3.75
for head_t, body_t in bullets:
    text(s, 1.25, by, 11.0, 0.5,
         [[("• ", 13, PURPLE, True), (head_t, 13, WHITE, True), (body_t, 13, MUTED, False)]])
    by += 0.52
footer(s, 7)

# ───────────────────────────────────────── Slide 8 — Cloud + Rule Distribution
s = slide()
header(s, "CAPABILITY 4 & 5", "Cloud Service + Rule Distribution")
# flow
chip(s, 0.85, 1.9, 2.3, 0.7, [[("CTI feeds", 12, GREEN, True)], [("OpenCTI / Avocado", 9, MUTED, False)]], GREEN)
arrow(s, 3.2, 2.16, 0.3, GREEN)
chip(s, 3.6, 1.9, 2.3, 0.7, [[("Scoring +", 12, GREEN, True)], [("Policy Engine", 11, MUTED, False)]], GREEN)
arrow(s, 5.95, 2.16, 0.3, GREEN)
chip(s, 6.35, 1.9, 2.9, 0.7, [[("Signed Policy Artifact", 12, PURPLE, True)], [("blacklist · STIX · IDS · TLS fp", 9, MUTED, False)]], PURPLE)
arrow(s, 9.3, 2.16, 0.3, GREEN)
chip(s, 9.7, 1.9, 2.75, 0.7, [[("Edge / NDR / FW", 12, WHITE, True)], [("HTTP · MQTT · SSE", 9, MUTED, False)]], WHITE)
text(s, 0.85, 2.95, 11.6, 0.4, [[("Two distribution pipelines:", 13, WHITE, True)]])
table(s, 0.85, 3.4, 11.6,
      ["Pipeline", "Direction", "Channel", "Content"],
      [["CTI / IoC blacklist", "cloud → edge / NDR / FW", "HTTP + MQTT + SSE", "IP / domain / hash · IDS rules · TLS fingerprints"],
       ["OT detection policy", "cloud → edge", "MQTT", "enabled rules + baselines"]],
      col_w=[2.7, 3.0, 2.4, 3.5], row_h=0.6, head_h=0.46, fs=11)
chip(s, 0.85, 5.55, 11.6, 0.95,
     [[("GAP → PQC + PUF", 14, PURPLE, True)],
      [("Distribution relies on a copyable API key + optional TLS. Artifacts are not end-to-end signed; "
        "channels are not quantum-resistant.", 12.5, WHITE, False)]], PURPLE)
footer(s, 8)

# ───────────────────────────────────────── Slide 9 — The Trust Gap
s = slide()
header(s, "THE PROBLEM", "The Trust Gap — Current State vs. Need")
table(s, 0.85, 1.78, 11.6,
      ["Dimension", "Today", "Gap", "Owner"],
      [["Device identity", "Copyable API key + invite", "No hardware-bound identity", ("PUF", PURPLE, True)],
       ["Behavioral trust", "Mature scoring engine", "Behavior-only; no identity/context", ("fuse", GREEN, True)],
       ["Channels", "TLS (ECDHE/RSA)", "Not quantum-safe", ("PQC", PURPLE, True)],
       ["Policy integrity", "API key, no signing", "Tamper / replay possible", ("PQC", PURPLE, True)],
       ["Access control", "Portal RBAC only", "No device-level zero trust", ("ZTNA", PURPLE, True)],
       ["Enforcement", "Recommendations only", "No isolation on trust drop", ("ZTNA", PURPLE, True)],
       ["Certificates", "Static, manual", "No HW root / rotation / PQC", ("PUF+PQC", PURPLE, True)]],
      col_w=[2.5, 3.2, 4.0, 1.9], row_h=0.5, head_h=0.44, fs=11)
chip(s, 0.85, 6.05, 11.6, 0.62,
     [[("We have behavioral trust. We need identity, channel, and access trust.", 14, GREEN, True)]], GREEN)
footer(s, 9)

# ───────────────────────────────────────── Slide 10 — Unified Trust Layer
s = slide()
header(s, "THE TARGET", "The Unified Trust Layer")
# stacked layers PUF->PQC->ZTNA + behavioral feeding in
defs = [
    ("PUF — Identity Root", "unclonable device identity · key wrapping · attestation", PURPLE),
    ("PQC — Crypto", "quantum-safe KEM & signatures · mTLS · artifact signing", PURPLE),
    ("ZTNA — Access", "device posture · dynamic authz · micro-segmentation", PURPLE),
]
y = 1.9
for i, (t, d, col) in enumerate(defs):
    box(s, 1.4, y, 6.6, 0.92, fill=BG2, line=col, line_w=1.25, round_=True)
    text(s, 1.65, y + 0.12, 6.1, 0.4, [[(t, 14, col, True)]])
    text(s, 1.65, y + 0.5, 6.1, 0.35, [[(d, 10.5, MUTED, False)]])
    if i < 2:
        down_arrow(s, 4.6, y + 0.93, 0.22, color=GREEN)
    y += 1.16
# behavioral trust feeding ZTNA
box(s, 8.6, 1.9, 4.0, 1.08, fill=BG2, line=GREEN, line_w=1.25, round_=True)
text(s, 8.85, 2.05, 3.5, 0.4, [[("Behavioral Trust (today)", 13, GREEN, True)]])
text(s, 8.85, 2.45, 3.5, 0.4, [[("entity state · alerts", 10.5, MUTED, False)]])
# unified score box
box(s, 1.4, 5.6, 11.2, 1.0, fill=PURPLE, round_=True)
text(s, 1.6, 5.72, 10.8, 0.8,
     [[("Identity (PUF) → Channel (PQC) → Access (ZTNA) → Behavior (existing)  =  ", 14, WHITE, True),
       ("ONE Unified Trust Score", 15, BG, True)]],
     anchor=MSO_ANCHOR.MIDDLE)
down_arrow(s, 6.9, 5.0, 0.5, color=GREEN)
footer(s, 10)

# ───────────────────────────────────────── Slide 11 — Integration Points
s = slide()
header(s, "HOW WE CONNECT", "Partner Integration Points")
table(s, 0.85, 1.78, 11.6,
      ["#", "Integration", "Extension needed", "Partner"],
      [["1", "Device registration", "Add PUF attestation to enrollment", ("PUF", PURPLE, True)],
       ["2", "Health reporting", "Carry attestation nonce / proof", ("PUF", PURPLE, True)],
       ["3", "Policy artifacts", "Add signature + algo + key id; verify at edge", ("PQC", PURPLE, True)],
       ["4", "Channels", "Hybrid PQC-TLS + mTLS north/south", ("PQC+PUF", PURPLE, True)],
       ["5", "Trust decision input", "Consume entity-state / alert stream as input", ("ZTNA", PURPLE, True)],
       ["6", "Micro-segmentation", "Map trust graph to ZTNA segments", ("ZTNA", PURPLE, True)],
       ["7", "Enforcement", "Edge adapter ↔ ZTNA gateway / NAC", ("ZTNA", PURPLE, True)],
       ["8", "Unified score", "Add identity + channel components", ("All", GREEN, True)]],
      col_w=[0.6, 3.0, 5.7, 2.3], row_h=0.46, head_h=0.42, fs=11)
footer(s, 11)

# ───────────────────────────────────────── Slide 12 — Roadmap
s = slide()
header(s, "THE PLAN", "Phased Roadmap")
phases = [
    ("P0", "Baseline (today)", "Behavioral trust + CTI distribution + passive detection live", MUTED),
    ("P1", "Identity root", "PUF device identity · attestation · anti-spoof enrollment", PURPLE),
    ("P2", "Quantum-safe", "Hybrid PQC-TLS · signed & verified policy artifacts", PURPLE),
    ("P3", "Zero-trust access", "Trust state → ZTNA decisions · micro-seg · enforcement", PURPLE),
    ("P4", "Unified trust", "One score fusing identity, channel, behavior, context", GREEN),
]
y = 1.9
for code, name, desc, col in phases:
    chip(s, 0.85, y, 1.1, 0.86, [[(code, 18, col, True)]], col)
    box(s, 2.1, y, 10.35, 0.86, fill=BG2, line=LINE, line_w=0.75, round_=True)
    text(s, 2.4, y + 0.12, 9.9, 0.4, [[(name, 14, col if col != MUTED else WHITE, True)]])
    text(s, 2.4, y + 0.5, 9.9, 0.35, [[(desc, 11.5, MUTED, False)]])
    y += 1.0
text(s, 0.85, 7.0, 11.6, 0.4,
     [[("Incremental · backward-compatible · reversible (hybrid before cutover, audit before enforce)",
        11, GREEN, True)]], align=PP_ALIGN.CENTER)
footer(s, 12)

# ───────────────────────────────────────── Slide 13 — Open Questions
s = slide()
header(s, "DISCUSSION", "Open Questions for Partners")
qs = [
    ("Latency", "Can PQC handshakes / signature checks meet OT control-loop timing on constrained edge hardware?"),
    ("PUF carrier", "Do target gateways have a PUF source (TPM / SE / SRAM-PUF), or is a secure element required?"),
    ("Offline behavior", "How to verify attestation / artifacts when the edge is disconnected (grace period? cached root)?"),
    ("Safety", "Fail-open vs. fail-secure for ZTNA isolation in a live OT environment?"),
    ("Key governance", "Rotation, revocation, and alignment with our tenant model."),
    ("Standards", "PQC algorithm selection (ML-KEM / ML-DSA) and migration timeline vs. NIST guidance."),
]
y = 1.85
for t, d in qs:
    box(s, 0.85, y, 0.12, 0.74, fill=GREEN)
    box(s, 1.05, y, 11.4, 0.74, fill=BG2, line=LINE, line_w=0.6, round_=True)
    text(s, 1.3, y + 0.08, 2.7, 0.6, [[(t, 13, GREEN, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.7, y + 0.06, 8.5, 0.62, [[(d, 11.5, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.86
footer(s, 13)

# ───────────────────────────────────────── Slide 14 — Closing
s = slide()
box(s, 0, 0, 13.333, 0.10, fill=GREEN)
box(s, 0, 7.40, 13.333, 0.10, fill=PURPLE)
lw = 4.6
lh = lw * 682 / 1024
s.shapes.add_picture(LOGO, Emu(int((13.333 - lw) / 2 * EMU)),
                     Emu(int(1.7 * EMU)), Emu(int(lw * EMU)), Emu(int(lh * EMU)))
text(s, 1.0, 4.7, 11.333, 0.7,
     [[("Let's build the OT Trust Layer together.", 26, WHITE, True)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 5.5, 11.333, 0.5,
     [[("Identity (PUF)  ·  Channel (PQC)  ·  Access (ZTNA)  ·  Behavior (SenseL)", 14, GREEN, True)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 6.4, 11.333, 0.4,
     [[("SenseL EdgeX  ·  by AvocadoAI", 12, MUTED, False)]], align=PP_ALIGN.CENTER)

prs.save(OUT)
print("Saved:", OUT, "slides:", len(prs.slides._sldIdLst))
